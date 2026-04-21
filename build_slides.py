"""
Build a polished PowerPoint deck for the COMP258 Final Project.

Usage:
    python build_slides.py

Produces:
    ./COMP258_FinalProject_Slides.pptx
    ./slides_assets/*.png   (charts embedded in the deck)
"""

import os
import joblib
import numpy as np
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---------------------------------------------------------------------------
# Paths
# ---------------------------------------------------------------------------
HERE = os.path.abspath(os.path.dirname(__file__))
ASSETS = os.path.join(HERE, "slides_assets")
os.makedirs(ASSETS, exist_ok=True)

METRICS_DIR = os.path.join(HERE, "backend", "saved_models")
OUT_PPTX = os.path.join(HERE, "COMP258_FinalProject_Slides.pptx")


# ---------------------------------------------------------------------------
# Theme
# ---------------------------------------------------------------------------
PRIMARY = RGBColor(0x1E, 0x3A, 0x8A)    # indigo
ACCENT  = RGBColor(0xF5, 0x9E, 0x0B)    # amber
TEAL    = RGBColor(0x0E, 0xA5, 0xE9)    # sky
GREEN   = RGBColor(0x10, 0xB9, 0x81)    # emerald
RED     = RGBColor(0xEF, 0x44, 0x44)    # red
DARK    = RGBColor(0x0F, 0x17, 0x2A)    # slate-900
MUTED   = RGBColor(0x64, 0x74, 0x8B)    # slate-500
LIGHT   = RGBColor(0xF8, 0xFA, 0xFC)    # slate-50
CARD    = RGBColor(0xE2, 0xE8, 0xF0)    # slate-200
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)

PRIMARY_HEX = "#1E3A8A"
ACCENT_HEX  = "#F59E0B"
TEAL_HEX    = "#0EA5E9"
GREEN_HEX   = "#10B981"
DARK_HEX    = "#0F172A"
MUTED_HEX   = "#64748B"
LIGHT_HEX   = "#F8FAFC"


# ---------------------------------------------------------------------------
# Load real model metrics
# ---------------------------------------------------------------------------
def load_metrics():
    out = {}
    for name in ("persistence", "gpa", "atrisk", "stacked"):
        path = os.path.join(METRICS_DIR, f"{name}_metrics.pkl")
        out[name] = joblib.load(path) if os.path.exists(path) else {}
    return out


METRICS = load_metrics()


# ---------------------------------------------------------------------------
# Chart helpers
# ---------------------------------------------------------------------------
plt.rcParams.update({
    "font.family": "DejaVu Sans",
    "font.size": 11,
    "axes.edgecolor": DARK_HEX,
    "axes.labelcolor": DARK_HEX,
    "xtick.color": DARK_HEX,
    "ytick.color": DARK_HEX,
    "axes.titleweight": "bold",
    "axes.titlecolor": DARK_HEX,
    "axes.spines.top": False,
    "axes.spines.right": False,
})


def chart_loss_curve(name, title, path):
    m = METRICS.get(name, {})
    hist = m.get("training_history") or {}
    loss = hist.get("loss") or []
    val_loss = hist.get("val_loss") or []
    if not loss:
        return None
    fig, ax = plt.subplots(figsize=(6.6, 3.6), dpi=150)
    ax.plot(loss, label="train loss", color=PRIMARY_HEX, linewidth=2.2)
    if val_loss:
        ax.plot(val_loss, label="val loss", color=ACCENT_HEX, linewidth=2.2)
    ax.set_title(title)
    ax.set_xlabel("Epoch")
    ax.set_ylabel("Loss")
    ax.grid(True, linestyle="--", alpha=0.35)
    ax.legend(frameon=False, loc="upper right")
    fig.tight_layout()
    fig.savefig(path, dpi=180, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    return path


def chart_metric_curve(name, title, metric_key, val_key, ylabel, path):
    m = METRICS.get(name, {})
    hist = m.get("training_history") or {}
    ys = hist.get(metric_key) or []
    vys = hist.get(val_key) or []
    if not ys:
        return None
    fig, ax = plt.subplots(figsize=(6.6, 3.6), dpi=150)
    ax.plot(ys, label=f"train {ylabel.lower()}", color=TEAL_HEX, linewidth=2.2)
    if vys:
        ax.plot(vys, label=f"val {ylabel.lower()}", color=GREEN_HEX, linewidth=2.2)
    ax.set_title(title)
    ax.set_xlabel("Epoch")
    ax.set_ylabel(ylabel)
    ax.grid(True, linestyle="--", alpha=0.35)
    ax.legend(frameon=False)
    fig.tight_layout()
    fig.savefig(path, dpi=180, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    return path


def chart_classifier_metrics(path):
    pm = METRICS.get("persistence", {})
    am = METRICS.get("atrisk", {})
    labels = ["Accuracy", "Precision", "Recall", "F1"]
    p_vals = [pm.get("accuracy", 0), pm.get("precision", 0),
              pm.get("recall", 0), pm.get("f1_score", 0)]
    a_vals = [am.get("accuracy", 0), am.get("precision", 0),
              am.get("recall", 0), am.get("f1_score", 0)]
    x = np.arange(len(labels))
    w = 0.38

    fig, ax = plt.subplots(figsize=(7.6, 3.8), dpi=150)
    b1 = ax.bar(x - w/2, p_vals, w, label="Persistence", color=PRIMARY_HEX)
    b2 = ax.bar(x + w/2, a_vals, w, label="At-Risk",     color=ACCENT_HEX)
    for bars in (b1, b2):
        for rect in bars:
            ax.text(rect.get_x() + rect.get_width()/2, rect.get_height() + 0.015,
                    f"{rect.get_height():.2f}", ha="center", va="bottom",
                    fontsize=9, color=DARK_HEX)
    ax.set_ylim(0, 1.05)
    ax.set_xticks(x); ax.set_xticklabels(labels)
    ax.set_title("Classifier evaluation metrics (test set)")
    ax.legend(frameon=False)
    ax.grid(True, axis="y", linestyle="--", alpha=0.35)
    fig.tight_layout()
    fig.savefig(path, dpi=180, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    return path


def chart_regressor_metrics(path):
    gm = METRICS.get("gpa", {})
    sm = METRICS.get("stacked", {})
    labels = ["R²", "MAE", "RMSE"]
    g_vals = [gm.get("r2_score", 0), gm.get("mae", 0), gm.get("rmse", 0)]
    s_vals = [sm.get("r2_score", 0), sm.get("mae", 0), sm.get("rmse", 0)]
    x = np.arange(len(labels)); w = 0.38

    fig, ax = plt.subplots(figsize=(7.6, 3.8), dpi=150)
    b1 = ax.bar(x - w/2, g_vals, w, label="GPA (stage 1)",     color=TEAL_HEX)
    b2 = ax.bar(x + w/2, s_vals, w, label="Stacked (stage 2)", color=GREEN_HEX)
    for bars in (b1, b2):
        for rect in bars:
            ax.text(rect.get_x() + rect.get_width()/2, rect.get_height() + 0.01,
                    f"{rect.get_height():.3f}", ha="center", va="bottom",
                    fontsize=9, color=DARK_HEX)
    ax.set_xticks(x); ax.set_xticklabels(labels)
    ax.set_title("Regression evaluation metrics (test set)")
    ax.legend(frameon=False)
    ax.grid(True, axis="y", linestyle="--", alpha=0.35)
    fig.tight_layout()
    fig.savefig(path, dpi=180, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    return path


def chart_confusion(name, title, path):
    m = METRICS.get(name, {})
    cm = m.get("confusion_matrix")
    if not cm:
        return None
    cm = np.array(cm)
    fig, ax = plt.subplots(figsize=(4.8, 3.8), dpi=150)
    im = ax.imshow(cm, cmap="Blues")
    for i in range(cm.shape[0]):
        for j in range(cm.shape[1]):
            ax.text(j, i, str(cm[i, j]),
                    ha="center", va="center", fontsize=14, fontweight="bold",
                    color="white" if cm[i, j] > cm.max()/2 else DARK_HEX)
    ax.set_xticks([0, 1]); ax.set_yticks([0, 1])
    ax.set_xticklabels(["Pred 0", "Pred 1"])
    ax.set_yticklabels(["True 0", "True 1"])
    ax.set_title(title)
    fig.colorbar(im, ax=ax, fraction=0.04, pad=0.04)
    fig.tight_layout()
    fig.savefig(path, dpi=180, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    return path


# Build all chart images
charts = {
    "persistence_loss": chart_loss_curve("persistence",
        "Persistence - training loss", os.path.join(ASSETS, "persistence_loss.png")),
    "gpa_loss":         chart_loss_curve("gpa",
        "GPA - training loss (MSE)", os.path.join(ASSETS, "gpa_loss.png")),
    "atrisk_loss":      chart_loss_curve("atrisk",
        "At-Risk - training loss", os.path.join(ASSETS, "atrisk_loss.png")),
    "stacked_loss":     chart_loss_curve("stacked",
        "Stacked Academic - training loss", os.path.join(ASSETS, "stacked_loss.png")),
    "clf_metrics":      chart_classifier_metrics(os.path.join(ASSETS, "clf_metrics.png")),
    "reg_metrics":      chart_regressor_metrics(os.path.join(ASSETS, "reg_metrics.png")),
    "persistence_cm":   chart_confusion("persistence",
        "Persistence - confusion matrix", os.path.join(ASSETS, "persistence_cm.png")),
    "atrisk_cm":        chart_confusion("atrisk",
        "At-Risk - confusion matrix", os.path.join(ASSETS, "atrisk_cm.png")),
}


# ---------------------------------------------------------------------------
# Slide primitives
# ---------------------------------------------------------------------------
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SLIDE_W, SLIDE_H = prs.slide_width, prs.slide_height

BLANK = prs.slide_layouts[6]


def add_rect(slide, x, y, w, h, fill, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
    shape.shadow.inherit = False
    return shape


def add_round_rect(slide, x, y, w, h, fill, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.adjustments[0] = 0.08
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
    shape.shadow.inherit = False
    return shape


def add_text(slide, x, y, w, h, text, *, size=18, bold=False,
             color=DARK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             font="Calibri"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor

    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        r.font.name = font
    return tb


def add_bullets(slide, x, y, w, h, items, *, size=16, color=DARK,
                bullet="•", bold=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(6)
        if isinstance(item, tuple):
            head, rest = item
            r = p.add_run()
            r.text = f"{bullet}  {head}"
            r.font.size = Pt(size); r.font.bold = True
            r.font.color.rgb = color; r.font.name = "Calibri"
            r2 = p.add_run()
            r2.text = f"  {rest}"
            r2.font.size = Pt(size); r2.font.bold = False
            r2.font.color.rgb = color; r2.font.name = "Calibri"
        else:
            r = p.add_run()
            r.text = f"{bullet}  {item}"
            r.font.size = Pt(size); r.font.bold = bold
            r.font.color.rgb = color; r.font.name = "Calibri"
    return tb


def add_page_footer(slide, page_num, total):
    # Thin accent bar at bottom
    add_rect(slide, 0, SLIDE_H - Inches(0.35), SLIDE_W, Inches(0.35), PRIMARY)
    add_rect(slide, 0, SLIDE_H - Inches(0.37), SLIDE_W, Inches(0.02), ACCENT)
    add_text(slide, Inches(0.5), SLIDE_H - Inches(0.33),
             Inches(8), Inches(0.3),
             "COMP258  ·  Student Success Prediction  ·  Full-Stack Intelligent App",
             size=11, color=WHITE)
    add_text(slide, SLIDE_W - Inches(1.3), SLIDE_H - Inches(0.33),
             Inches(0.9), Inches(0.3), f"{page_num} / {total}",
             size=11, color=WHITE, align=PP_ALIGN.RIGHT)


def add_slide_title_band(slide, title, subtitle=None):
    # Top ribbon
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.15), ACCENT)
    # Title text
    add_text(slide, Inches(0.5), Inches(0.28),
             SLIDE_W - Inches(1.0), Inches(0.7),
             title, size=30, bold=True, color=PRIMARY)
    if subtitle:
        add_text(slide, Inches(0.5), Inches(0.95),
                 SLIDE_W - Inches(1.0), Inches(0.45),
                 subtitle, size=15, color=MUTED)
    # Underline accent
    add_rect(slide, Inches(0.5), Inches(1.45),
             Inches(1.2), Inches(0.06), ACCENT)


# Count slides first so footers can show "x / total"
SLIDES = []  # list of callables — we render after counting

def slide(fn):
    SLIDES.append(fn)
    return fn


# ---------------------------------------------------------------------------
# Slide 1: Title
# ---------------------------------------------------------------------------
@slide
def s_title(s):
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, PRIMARY)
    # Decorative accent block
    add_rect(s, 0, Inches(5.6), SLIDE_W, Inches(0.08), ACCENT)
    add_rect(s, Inches(10.5), 0, Inches(3.0), SLIDE_H, RGBColor(0x17, 0x2B, 0x6E))
    add_rect(s, Inches(10.5), 0, Inches(0.08), SLIDE_H, ACCENT)

    add_text(s, Inches(0.8), Inches(1.0), Inches(9.5), Inches(0.5),
             "COMP258  ·  Neural Networks  ·  Final Project",
             size=16, color=ACCENT, bold=True)
    add_text(s, Inches(0.8), Inches(1.6), Inches(10), Inches(1.8),
             "Student Success Prediction",
             size=54, bold=True, color=WHITE)
    add_text(s, Inches(0.8), Inches(3.2), Inches(10), Inches(1.2),
             "A full-stack intelligent web app that forecasts\npersistence, GPA, and at-risk status using deep learning.",
             size=22, color=RGBColor(0xCB, 0xD5, 0xE1))

    add_text(s, Inches(0.8), Inches(5.85), Inches(9), Inches(0.45),
             "Group:  [YourGroupName]",
             size=18, bold=True, color=WHITE)
    add_text(s, Inches(0.8), Inches(6.3), Inches(9), Inches(0.45),
             "Members:  [Name 1] · [Name 2] · [Name 3] · [Name 4] · [Name 5]",
             size=14, color=RGBColor(0xCB, 0xD5, 0xE1))
    add_text(s, Inches(0.8), Inches(6.75), Inches(9), Inches(0.4),
             "Presented in Week 14  ·  Centennial College",
             size=13, color=RGBColor(0xCB, 0xD5, 0xE1))


# ---------------------------------------------------------------------------
# Slide 2: Agenda
# ---------------------------------------------------------------------------
@slide
def s_agenda(s):
    add_slide_title_band(s, "Agenda", "What we will cover today")

    items = [
        ("01", "Problem & Motivation"),
        ("02", "Dataset & Preprocessing"),
        ("03", "System Architecture"),
        ("04", "Neural Network Models (×4)"),
        ("05", "Training Strategy"),
        ("06", "Results & Evaluation"),
        ("07", "Live Demo"),
        ("08", "Challenges & Future Work"),
    ]
    x0, y0 = Inches(0.6), Inches(1.9)
    card_w, card_h = Inches(6.0), Inches(0.6)
    for i, (num, label) in enumerate(items):
        col = i % 2
        row = i // 2
        x = x0 + col * (card_w + Inches(0.25))
        y = y0 + row * (card_h + Inches(0.2))
        add_round_rect(s, x, y, card_w, card_h, LIGHT)
        add_rect(s, x, y, Inches(0.08), card_h, ACCENT)
        add_text(s, x + Inches(0.25), y, Inches(0.8), card_h,
                 num, size=18, bold=True, color=PRIMARY,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + Inches(1.0), y, card_w - Inches(1.1), card_h,
                 label, size=18, color=DARK, anchor=MSO_ANCHOR.MIDDLE)


# ---------------------------------------------------------------------------
# Slide 3: Problem & Motivation
# ---------------------------------------------------------------------------
@slide
def s_problem(s):
    add_slide_title_band(s, "Problem & Motivation",
                         "Why predicting student outcomes matters")
    # Left column - narrative
    add_text(s, Inches(0.6), Inches(1.9), Inches(6.4), Inches(0.5),
             "The challenge", size=18, bold=True, color=PRIMARY)
    add_bullets(s, Inches(0.6), Inches(2.4), Inches(6.4), Inches(4.5), [
        "Post-secondary institutions lose a meaningful share of first-year students.",
        "Early identification of at-risk students enables timely intervention.",
        "Instructors need lightweight, data-driven tools, not black-box dashboards.",
        "Client asks for an admin-friendly app that turns existing data into action.",
    ], size=15)

    # Right column - stat cards
    cards = [
        ("4", "neural networks trained", PRIMARY),
        ("~1,400", "student records", TEAL),
        ("14-15", "features per model", GREEN),
        ("REST + React", "full-stack delivery", ACCENT),
    ]
    cx, cy = Inches(7.3), Inches(1.9)
    cw, ch = Inches(2.75), Inches(1.25)
    for i, (big, small, col) in enumerate(cards):
        r = i // 2; c = i % 2
        x = cx + c * (cw + Inches(0.2))
        y = cy + r * (ch + Inches(0.2))
        add_round_rect(s, x, y, cw, ch, LIGHT)
        add_rect(s, x, y, cw, Inches(0.08), col)
        add_text(s, x, y + Inches(0.18), cw, Inches(0.6),
                 big, size=28, bold=True, color=col, align=PP_ALIGN.CENTER)
        add_text(s, x, y + Inches(0.78), cw, Inches(0.45),
                 small, size=12, color=MUTED, align=PP_ALIGN.CENTER)


# ---------------------------------------------------------------------------
# Slide 4: Objectives
# ---------------------------------------------------------------------------
@slide
def s_objectives(s):
    add_slide_title_band(s, "Project Objectives",
                         "Four prediction tasks delivered as one intelligent app")

    tasks = [
        ("01", "First-Year Persistence",
         "Binary classification: will this student continue after year 1?", PRIMARY),
        ("02", "Second-Term GPA",
         "Regression from 1st-term GPA, HS marks and demographics.", TEAL),
        ("03", "At-Risk Detection",
         "Flag students needing support (persistence = 0 OR GPA < 2.0).", RED),
        ("04", "Stacked Academic Score",
         "Two-stage regression for an overall academic-performance score.", GREEN),
    ]
    x0, y0 = Inches(0.55), Inches(1.9)
    cw, ch = Inches(6.1), Inches(2.35)
    for i, (num, title, desc, col) in enumerate(tasks):
        r = i // 2; c = i % 2
        x = x0 + c * (cw + Inches(0.25))
        y = y0 + r * (ch + Inches(0.2))
        add_round_rect(s, x, y, cw, ch, LIGHT)
        add_rect(s, x, y, Inches(0.12), ch, col)
        add_text(s, x + Inches(0.35), y + Inches(0.15),
                 Inches(1.1), Inches(0.6), num,
                 size=34, bold=True, color=col)
        add_text(s, x + Inches(1.5), y + Inches(0.2),
                 cw - Inches(1.7), Inches(0.7), title,
                 size=20, bold=True, color=DARK)
        add_text(s, x + Inches(1.5), y + Inches(0.9),
                 cw - Inches(1.7), ch - Inches(1.0), desc,
                 size=14, color=MUTED)


# ---------------------------------------------------------------------------
# Slide 5: Dataset Overview
# ---------------------------------------------------------------------------
@slide
def s_dataset(s):
    add_slide_title_band(s, "Dataset Overview",
                         "Centennial College student records (NDA-protected)")
    add_text(s, Inches(0.6), Inches(1.9), Inches(7.5), Inches(0.5),
             "Input features (14-15 per model)", size=18, bold=True, color=PRIMARY)
    add_bullets(s, Inches(0.6), Inches(2.4), Inches(7.5), Inches(4.5), [
        ("Academic:", "first / second term GPA, HS average, math score, English grade"),
        ("Program:", "school, funding, fast-track, coop, residency, prev. education"),
        ("Demographic:", "gender, age group, first language"),
        ("Target labels:", "first-year persistence (0/1), GPA (0.0-4.5), derived at-risk"),
    ], size=14)

    # Right: small "data card"
    x = Inches(8.3); y = Inches(1.9); w = Inches(4.5); h = Inches(4.3)
    add_round_rect(s, x, y, w, h, LIGHT)
    add_rect(s, x, y, w, Inches(0.08), TEAL)
    add_text(s, x, y + Inches(0.2), w, Inches(0.45),
             "Dataset snapshot", size=17, bold=True, color=PRIMARY,
             align=PP_ALIGN.CENTER)
    facts = [
        ("~1,400", "records"),
        ("15", "raw variables"),
        ("Numeric + Categorical", "mixed types"),
        ("Imputed", "median / mode per column"),
        ("Standardized", "StandardScaler on features"),
        ("80 / 20", "train / test split"),
    ]
    row_y = y + Inches(0.85)
    for k, v in facts:
        add_text(s, x + Inches(0.3), row_y, Inches(2.0), Inches(0.35),
                 k, size=14, bold=True, color=DARK)
        add_text(s, x + Inches(2.2), row_y, Inches(2.2), Inches(0.35),
                 v, size=13, color=MUTED, align=PP_ALIGN.RIGHT)
        row_y += Inches(0.5)


# ---------------------------------------------------------------------------
# Slide 6: Data pipeline
# ---------------------------------------------------------------------------
@slide
def s_pipeline(s):
    add_slide_title_band(s, "Data Preparation Pipeline",
                         "From raw CSV to model-ready tensors")

    steps = [
        ("1", "Ingest", "Load CSV →\npersist to SQLite", TEAL),
        ("2", "Clean", "Drop rows w/o target\nImpute features", PRIMARY),
        ("3", "Split", "80 / 20\nstratified where possible", ACCENT),
        ("4", "Scale", "StandardScaler\nfit on train only", GREEN),
        ("5", "Train", "Keras Sequential\nearly-stopping + LR decay", RED),
    ]
    x = Inches(0.6); y = Inches(2.2)
    w = Inches(2.35); h = Inches(2.3)
    gap = Inches(0.12)
    total_w = 5 * w + 4 * gap
    x_start = (SLIDE_W - total_w) / 2
    for i, (num, title, body, col) in enumerate(steps):
        sx = x_start + i * (w + gap)
        add_round_rect(s, sx, y, w, h, LIGHT)
        add_rect(s, sx, y, w, Inches(0.08), col)
        add_text(s, sx, y + Inches(0.2), w, Inches(0.55),
                 num, size=32, bold=True, color=col, align=PP_ALIGN.CENTER)
        add_text(s, sx, y + Inches(0.85), w, Inches(0.5),
                 title, size=18, bold=True, color=DARK, align=PP_ALIGN.CENTER)
        add_text(s, sx + Inches(0.15), y + Inches(1.35),
                 w - Inches(0.3), Inches(0.9),
                 body, size=12, color=MUTED, align=PP_ALIGN.CENTER)
        # arrow
        if i < 4:
            add_text(s, sx + w - Inches(0.05), y + h / 2 - Inches(0.2),
                     gap + Inches(0.2), Inches(0.45),
                     "→", size=26, bold=True, color=ACCENT,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    add_text(s, Inches(0.6), Inches(5.0), Inches(12), Inches(0.5),
             "Why imputation?",
             size=16, bold=True, color=PRIMARY)
    add_bullets(s, Inches(0.6), Inches(5.4), Inches(12), Inches(1.5), [
        "Dropping any missing value collapsed the dataset from ~1,400 rows to ~500 — too little signal.",
        "We only drop rows whose target label is missing, then impute feature NaNs (median for numeric, mode for categorical).",
        "This preserves variance while keeping labels honest.",
    ], size=13)


# ---------------------------------------------------------------------------
# Slide 7: System architecture
# ---------------------------------------------------------------------------
@slide
def s_architecture(s):
    add_slide_title_band(s, "System Architecture",
                         "React SPA  →  Flask REST API  →  TensorFlow models + SQLite")

    y = Inches(2.1); h = Inches(3.6)
    tiles = [
        ("Frontend", "React 18 SPA", [
            "Dashboard, Data Explorer",
            "Prediction forms for each task",
            "Model-metrics panel, Admin panel",
            "Axios client · Chart.js",
        ], PRIMARY),
        ("Backend", "Flask + Blueprints", [
            "MVC: routes / services / models",
            "/api/data · /api/predict",
            "/api/models · /api/admin",
            "CORS enabled for SPA",
        ], TEAL),
        ("ML Layer", "TensorFlow / Keras", [
            "4 Sequential neural nets",
            "StandardScaler per model",
            ".keras + .pkl artefacts",
            "Retrain endpoint",
        ], ACCENT),
        ("Storage", "SQLite + CSV seed", [
            "SQLAlchemy ORM",
            "Student + PredictionLog tables",
            "CSV re-import endpoint",
            "Docker-ready (compose)",
        ], GREEN),
    ]
    col_w = Inches(3.0); gap = Inches(0.15)
    total = 4 * col_w + 3 * gap
    x_start = (SLIDE_W - total) / 2
    for i, (head, sub, items, col) in enumerate(tiles):
        x = x_start + i * (col_w + gap)
        add_round_rect(s, x, y, col_w, h, LIGHT)
        add_rect(s, x, y, col_w, Inches(0.5), col)
        add_text(s, x, y, col_w, Inches(0.5), head,
                 size=16, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + Inches(0.2), y + Inches(0.6),
                 col_w - Inches(0.4), Inches(0.45),
                 sub, size=13, bold=True, color=DARK)
        bul = add_bullets(s, x + Inches(0.2), y + Inches(1.05),
                          col_w - Inches(0.4), h - Inches(1.1),
                          items, size=11, color=MUTED)
    add_text(s, 0, Inches(5.95), SLIDE_W, Inches(0.4),
             "HTTP / JSON  ·  RESTful endpoints  ·  stateless backend",
             size=13, color=MUTED, align=PP_ALIGN.CENTER)


# ---------------------------------------------------------------------------
# Slide 8: Backend details
# ---------------------------------------------------------------------------
@slide
def s_backend(s):
    add_slide_title_band(s, "Backend  ·  Flask REST API",
                         "MVC blueprints with clear separation of concerns")

    # Left: endpoint table
    add_text(s, Inches(0.6), Inches(1.9), Inches(7), Inches(0.4),
             "Key endpoints", size=17, bold=True, color=PRIMARY)
    endpoints = [
        ("GET",  "/api/data/summary",        "dataset stats"),
        ("GET",  "/api/data/students",       "paginated + filtered"),
        ("GET",  "/api/data/correlation",    "feature correlation matrix"),
        ("POST", "/api/predict/persistence", "classify persistence"),
        ("POST", "/api/predict/gpa",         "regress 2nd-term GPA"),
        ("POST", "/api/predict/atrisk",      "classify risk level"),
        ("POST", "/api/predict/pipeline",    "two-stage academic score"),
        ("GET",  "/api/models/metrics",      "live evaluation metrics"),
        ("POST", "/api/admin/retrain",       "re-train all four models"),
    ]
    header_y = Inches(2.4)
    row_h = Inches(0.38)
    for i, (method, path, desc) in enumerate(endpoints):
        rx, ry = Inches(0.6), header_y + i * row_h
        if i % 2 == 0:
            add_rect(s, rx, ry, Inches(7.6), row_h, LIGHT)
        col = TEAL if method == "GET" else ACCENT
        add_rect(s, rx + Inches(0.05), ry + Inches(0.06),
                 Inches(0.55), Inches(0.26), col)
        add_text(s, rx + Inches(0.05), ry + Inches(0.04),
                 Inches(0.55), Inches(0.3), method,
                 size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, rx + Inches(0.75), ry,
                 Inches(3.8), row_h, path,
                 size=12, bold=True, color=DARK,
                 anchor=MSO_ANCHOR.MIDDLE, font="Consolas")
        add_text(s, rx + Inches(4.55), ry,
                 Inches(3.1), row_h, desc,
                 size=11, color=MUTED, anchor=MSO_ANCHOR.MIDDLE)

    # Right: design patterns
    add_text(s, Inches(8.6), Inches(1.9), Inches(4.3), Inches(0.4),
             "Design patterns", size=17, bold=True, color=PRIMARY)
    add_bullets(s, Inches(8.6), Inches(2.3), Inches(4.5), Inches(3.5), [
        ("MVC:",        "routes · services · models"),
        ("Blueprint:",  "one per concern"),
        ("Repository:", "DB access via SQLAlchemy"),
        ("DI:",         "prediction_service injected at startup"),
        ("Factory:",    "create_app() builds the Flask app"),
    ], size=13)


# ---------------------------------------------------------------------------
# Slide 9: Frontend details
# ---------------------------------------------------------------------------
@slide
def s_frontend(s):
    add_slide_title_band(s, "Frontend  ·  React SPA",
                         "Component-driven, feature-organised")

    components = [
        ("Dashboard",            "home + KPIs", PRIMARY),
        ("DataExplorer",         "filter students", TEAL),
        ("PersistencePrediction","task form", ACCENT),
        ("GPAPrediction",        "task form", ACCENT),
        ("AtRiskPrediction",     "task form", ACCENT),
        ("AcademicPipeline",     "two-stage demo", GREEN),
        ("ModelMetrics",         "live metrics", RED),
        ("AdminPanel",           "retrain / reimport", PRIMARY),
    ]
    x0, y0 = Inches(0.55), Inches(2.0)
    cw, ch = Inches(3.0), Inches(0.95)
    gap_x, gap_y = Inches(0.2), Inches(0.2)
    for i, (name, desc, col) in enumerate(components):
        r = i // 4; c = i % 4
        x = x0 + c * (cw + gap_x); y = y0 + r * (ch + gap_y)
        add_round_rect(s, x, y, cw, ch, LIGHT)
        add_rect(s, x, y, Inches(0.1), ch, col)
        add_text(s, x + Inches(0.25), y + Inches(0.1), cw - Inches(0.3), Inches(0.45),
                 name, size=15, bold=True, color=DARK)
        add_text(s, x + Inches(0.25), y + Inches(0.5), cw - Inches(0.3), Inches(0.4),
                 desc, size=12, color=MUTED)

    add_text(s, Inches(0.55), Inches(4.55), Inches(12), Inches(0.4),
             "Key frontend choices",
             size=17, bold=True, color=PRIMARY)
    add_bullets(s, Inches(0.55), Inches(4.95), Inches(12), Inches(2), [
        "Single-Page App with react-router — no reloads between prediction screens.",
        "axios service module centralises API calls (src/services/api.js).",
        "Form inputs map directly to the backend feature contract for zero guesswork.",
        "Results cards reuse the same layout across all four models.",
    ], size=13)


# ---------------------------------------------------------------------------
# Slide 10: Common NN recipe
# ---------------------------------------------------------------------------
@slide
def s_nn_recipe(s):
    add_slide_title_band(s, "Neural Network Design — Shared Recipe",
                         "One proven architecture, four outputs")

    # Diagram of layers
    layers = [
        ("Input", "14-15 features", LIGHT, DARK),
        ("Dense 128 + BN + Drop 0.3", "ReLU · L2(0.001)", PRIMARY, WHITE),
        ("Dense 64  + BN + Drop 0.3", "ReLU · L2(0.001)", TEAL, WHITE),
        ("Dense 32  + Drop 0.2",      "ReLU",             GREEN, WHITE),
        ("Dense 1",                   "sigmoid / linear", ACCENT, WHITE),
    ]
    y = Inches(2.1); h = Inches(0.75); gap = Inches(0.12)
    x = Inches(0.8); w = Inches(5.8)
    for i, (head, sub, fill, txt) in enumerate(layers):
        ly = y + i * (h + gap)
        add_round_rect(s, x, ly, w, h, fill)
        add_text(s, x + Inches(0.25), ly, Inches(3.5), h, head,
                 size=14, bold=True, color=txt, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + Inches(3.6), ly, w - Inches(3.7), h, sub,
                 size=12, color=txt, anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.RIGHT)

    # Right: why this shape
    add_text(s, Inches(7.2), Inches(2.0), Inches(5.5), Inches(0.4),
             "Why this topology?", size=17, bold=True, color=PRIMARY)
    add_bullets(s, Inches(7.2), Inches(2.4), Inches(5.7), Inches(4.5), [
        ("Funnel widths:", "128 → 64 → 32 compresses features gradually."),
        ("BatchNorm:",     "stabilises training, speeds convergence."),
        ("Dropout:",       "0.3 / 0.2 / 0.1 fights overfitting on ~1k rows."),
        ("L2 (1e-3):",     "gentle weight decay on the two large layers."),
        ("Head:",          "sigmoid for classifiers, linear for regressors."),
        ("Optimizer:",     "Adam @ 1e-3 · reduced on plateau."),
    ], size=13)


# ---------------------------------------------------------------------------
# Slide 11-14: per-model
# ---------------------------------------------------------------------------
def make_model_slide(title, subtitle, key, chart_path, is_classifier):
    def _draw(s):
        add_slide_title_band(s, title, subtitle)

        m = METRICS.get(key, {})
        # Left: metric cards
        if is_classifier:
            values = [
                ("Accuracy",  m.get("accuracy", 0),  PRIMARY),
                ("Precision", m.get("precision", 0), TEAL),
                ("Recall",    m.get("recall", 0),    GREEN),
                ("F1",        m.get("f1_score", 0),  ACCENT),
            ]
            fmt = lambda v: f"{v*100:0.1f}%"
        else:
            values = [
                ("R²",   m.get("r2_score", 0), PRIMARY),
                ("MAE",  m.get("mae", 0),      TEAL),
                ("RMSE", m.get("rmse", 0),     GREEN),
                ("Epochs", m.get("epochs_trained", 0), ACCENT),
            ]
            def fmt_f(k, v):
                if k == "Epochs": return str(int(v))
                return f"{v:0.3f}"
            fmt = None  # handled below

        x0, y0 = Inches(0.55), Inches(1.95)
        cw, ch = Inches(2.55), Inches(1.1)
        for i, tpl in enumerate(values):
            k, v = tpl[0], tpl[1]; col = tpl[2]
            r = i // 2; c = i % 2
            x = x0 + c * (cw + Inches(0.15))
            y = y0 + r * (ch + Inches(0.15))
            add_round_rect(s, x, y, cw, ch, LIGHT)
            add_rect(s, x, y, cw, Inches(0.08), col)
            add_text(s, x, y + Inches(0.15), cw, Inches(0.45),
                     k, size=13, bold=True, color=MUTED, align=PP_ALIGN.CENTER)
            if is_classifier:
                disp = fmt(v)
            else:
                disp = fmt_f(k, v)
            add_text(s, x, y + Inches(0.55), cw, Inches(0.55),
                     disp, size=24, bold=True, color=col, align=PP_ALIGN.CENTER)

        # Samples line
        ts = m.get("training_samples", 0); te = m.get("test_samples", 0)
        ep = m.get("epochs_trained", 0)
        add_text(s, Inches(0.55), Inches(4.45),
                 Inches(5.4), Inches(0.4),
                 f"Train / Test: {ts} / {te} samples  ·  Epochs trained: {ep}",
                 size=12, color=MUTED)

        # Feature list summary
        add_text(s, Inches(0.55), Inches(4.95),
                 Inches(5.4), Inches(0.4),
                 "Inputs", size=15, bold=True, color=PRIMARY)
        input_map = {
            "persistence":
                "1st & 2nd term GPA · HS avg · math · English · school · funding · residency · age · gender · first language · fast-track · coop · prev. education",
            "gpa":
                "1st term GPA · HS avg · math · English · school · funding · residency · age · gender · first language · fast-track · coop · prev. education",
            "atrisk":
                "1st term GPA · HS avg · math · English · school · funding · residency · age · gender · first language · fast-track · coop · prev. education",
            "stacked":
                "1st + 2nd term GPA · HS avg · math · English · plus the full demographic set (14 features total)",
        }
        add_text(s, Inches(0.55), Inches(5.3),
                 Inches(5.8), Inches(1.8),
                 input_map.get(key, ""), size=11, color=MUTED)

        # Right: chart
        if chart_path and os.path.exists(chart_path):
            s.shapes.add_picture(chart_path, Inches(6.5), Inches(1.9),
                                 width=Inches(6.5), height=Inches(3.8))

        # If classifier, also drop the confusion matrix at bottom right
        cm_path = os.path.join(ASSETS, f"{key}_cm.png")
        if is_classifier and os.path.exists(cm_path):
            s.shapes.add_picture(cm_path, Inches(9.0), Inches(5.0),
                                 width=Inches(3.5), height=Inches(2.0))
    _draw.__name__ = f"s_model_{key}"
    return _draw

SLIDES.append(make_model_slide(
    "Model 1  ·  Persistence Classifier",
    "Binary classification: will the student continue past year 1?",
    "persistence", charts["persistence_loss"], is_classifier=True))

SLIDES.append(make_model_slide(
    "Model 2  ·  GPA Regressor",
    "Regression: predict second-term GPA (0.0 - 4.5)",
    "gpa", charts["gpa_loss"], is_classifier=False))

SLIDES.append(make_model_slide(
    "Model 3  ·  At-Risk Classifier",
    "Flag students with persistence = 0 OR first-term GPA < 2.0",
    "atrisk", charts["atrisk_loss"], is_classifier=True))

SLIDES.append(make_model_slide(
    "Model 4  ·  Stacked Academic Pipeline",
    "Stage-2 regression: GPAs + HS + demographics → overall score",
    "stacked", charts["stacked_loss"], is_classifier=False))


# ---------------------------------------------------------------------------
# Slide 15: Training strategy
# ---------------------------------------------------------------------------
@slide
def s_training(s):
    add_slide_title_band(s, "Training Strategy",
                         "Robust training without overfitting the small dataset")

    tactics = [
        ("EarlyStopping", "patience 15-20 on val_loss, restore best weights", PRIMARY),
        ("ReduceLROnPlateau", "halve LR if val_loss stalls 5-7 epochs", TEAL),
        ("Class weights", "balance persistence/at-risk imbalance", RED),
        ("StandardScaler", "fit on train only, applied consistently at inference", GREEN),
        ("Seeded RNGs", "Python · NumPy · TensorFlow (SEED = 42) → reproducible runs", ACCENT),
        ("Dropout + L2", "regularise the 128/64 layers against tiny-dataset overfit", PRIMARY),
    ]
    y0 = Inches(2.0); h = Inches(0.65); gap = Inches(0.12)
    for i, (name, desc, col) in enumerate(tactics):
        y = y0 + i * (h + gap)
        add_round_rect(s, Inches(0.6), y, Inches(12.1), h, LIGHT)
        add_rect(s, Inches(0.6), y, Inches(0.12), h, col)
        add_text(s, Inches(0.95), y, Inches(3.5), h,
                 name, size=15, bold=True, color=DARK, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(4.5), y, Inches(8.1), h,
                 desc, size=13, color=MUTED, anchor=MSO_ANCHOR.MIDDLE)


# ---------------------------------------------------------------------------
# Slide 16: Results summary
# ---------------------------------------------------------------------------
@slide
def s_results(s):
    add_slide_title_band(s, "Results  ·  Head-to-Head",
                         "Evaluation on the held-out 20 % test split")

    # Left: classifier chart
    if charts["clf_metrics"]:
        s.shapes.add_picture(charts["clf_metrics"],
                             Inches(0.5), Inches(1.95),
                             width=Inches(6.3), height=Inches(3.2))
    if charts["reg_metrics"]:
        s.shapes.add_picture(charts["reg_metrics"],
                             Inches(7.0), Inches(1.95),
                             width=Inches(6.3), height=Inches(3.2))

    # Bottom: key numbers
    pm = METRICS.get("persistence", {})
    am = METRICS.get("atrisk", {})
    gm = METRICS.get("gpa", {})
    sm = METRICS.get("stacked", {})

    def stat_card(x, y, label, value, color):
        w = Inches(2.95); h = Inches(1.35)
        add_round_rect(s, x, y, w, h, LIGHT)
        add_rect(s, x, y, w, Inches(0.08), color)
        add_text(s, x, y + Inches(0.15), w, Inches(0.45),
                 label, size=13, bold=True, color=MUTED, align=PP_ALIGN.CENTER)
        add_text(s, x, y + Inches(0.6), w, Inches(0.7),
                 value, size=26, bold=True, color=color, align=PP_ALIGN.CENTER)

    y = Inches(5.35)
    stat_card(Inches(0.45), y, "Persistence · F1",   f"{pm.get('f1_score', 0)*100:0.1f}%", PRIMARY)
    stat_card(Inches(3.55), y, "At-Risk · Accuracy", f"{am.get('accuracy', 0)*100:0.1f}%", RED)
    stat_card(Inches(6.65), y, "GPA · R²",           f"{gm.get('r2_score', 0):0.3f}",       TEAL)
    stat_card(Inches(9.75), y, "Stacked · R²",       f"{sm.get('r2_score', 0):0.3f}",       GREEN)


# ---------------------------------------------------------------------------
# Slide 17: Live demo
# ---------------------------------------------------------------------------
@slide
def s_demo(s):
    add_slide_title_band(s, "Live Demo",
                         "React SPA talking to the Flask REST API in real time")

    steps = [
        ("01", "Data Explorer", "Filter by school/gender/persistence — hits /api/data/students."),
        ("02", "Predict Persistence", "Fill the form → POST /api/predict/persistence → probability + label."),
        ("03", "Predict GPA", "Provide 1st-term GPA + HS scores → regression output."),
        ("04", "At-Risk Flag", "Three-tier: Low / Medium / High risk."),
        ("05", "Two-Stage Pipeline", "Stage 1 predicts 2nd-term GPA → fed into Stage 2 score."),
        ("06", "Admin / Retrain", "POST /api/admin/retrain re-fits all 4 models end-to-end."),
    ]
    x0, y0 = Inches(0.6), Inches(1.95)
    cw, ch = Inches(6.1), Inches(0.85)
    for i, (num, head, desc) in enumerate(steps):
        r = i // 2; c = i % 2
        x = x0 + c * (cw + Inches(0.2))
        y = y0 + r * (ch + Inches(0.15))
        add_round_rect(s, x, y, cw, ch, LIGHT)
        add_rect(s, x, y, Inches(0.9), ch, PRIMARY)
        add_text(s, x, y, Inches(0.9), ch, num,
                 size=20, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + Inches(1.05), y + Inches(0.1),
                 cw - Inches(1.1), Inches(0.35),
                 head, size=14, bold=True, color=DARK)
        add_text(s, x + Inches(1.05), y + Inches(0.42),
                 cw - Inches(1.1), Inches(0.5),
                 desc, size=11, color=MUTED)

    add_rect(s, Inches(0.6), Inches(5.8), Inches(12.1), Inches(0.9), LIGHT)
    add_rect(s, Inches(0.6), Inches(5.8), Inches(0.12), Inches(0.9), ACCENT)
    add_text(s, Inches(0.9), Inches(5.8), Inches(12), Inches(0.9),
             "💡 Tip for presenters: bring a copy of saved_models/ so the demo\n"
             "     works even if the venue has no internet. Start backend first, then frontend.",
             size=12, color=DARK, anchor=MSO_ANCHOR.MIDDLE)


# ---------------------------------------------------------------------------
# Slide 18: Challenges & Solutions
# ---------------------------------------------------------------------------
@slide
def s_challenges(s):
    add_slide_title_band(s, "Challenges & How We Solved Them",
                         "Small data, imbalanced labels, full-stack coordination")

    rows = [
        ("Missing values dropped the dataset in half",
         "Impute per-column (median / mode), keep all 1,400 rows.", PRIMARY),
        ("Class imbalance on persistence / at-risk",
         "Compute class_weight on the fly during training.", RED),
        ("Overfitting a small dataset",
         "L2, dropout, BatchNorm, EarlyStopping + ReduceLROnPlateau.", TEAL),
        ("Reproducible results for the demo",
         "Seed Python / NumPy / TensorFlow, restore best weights.", GREEN),
        ("Front-end ↔ back-end contract drift",
         "Single schema for each prediction form, mirrored in services.", ACCENT),
    ]
    y = Inches(2.0); h = Inches(0.82); gap = Inches(0.12)
    for i, (problem, solution, col) in enumerate(rows):
        yy = y + i * (h + gap)
        add_round_rect(s, Inches(0.55), yy, Inches(12.2), h, LIGHT)
        add_rect(s, Inches(0.55), yy, Inches(0.12), h, col)
        add_text(s, Inches(0.8), yy + Inches(0.08),
                 Inches(5.8), Inches(0.35),
                 "Challenge", size=10, bold=True, color=col)
        add_text(s, Inches(0.8), yy + Inches(0.38),
                 Inches(5.8), Inches(0.45),
                 problem, size=13, bold=True, color=DARK)
        add_text(s, Inches(6.8), yy + Inches(0.08),
                 Inches(5.8), Inches(0.35),
                 "Solution", size=10, bold=True, color=MUTED)
        add_text(s, Inches(6.8), yy + Inches(0.38),
                 Inches(5.8), Inches(0.45),
                 solution, size=13, color=DARK)


# ---------------------------------------------------------------------------
# Slide 19: Future Work
# ---------------------------------------------------------------------------
@slide
def s_future(s):
    add_slide_title_band(s, "Future Work",
                         "Where we would take the app next")
    ideas = [
        ("Model upgrades",
         "Try TabTransformer / gradient-boosted trees for tabular baseline parity.", PRIMARY),
        ("Explainability",
         "Per-prediction SHAP values in the UI so advisors see *why*.", ACCENT),
        ("Confidence calibration",
         "Temperature scaling for the two classifiers.", TEAL),
        ("Active learning loop",
         "Log low-confidence cases → advisor labels them → retrain nightly.", GREEN),
        ("Auth & auditing",
         "Role-based access for advisor vs. admin; immutable prediction log.", RED),
        ("Deployment",
         "Docker Compose → cloud (Azure/GCP), HTTPS, CI/CD.", PRIMARY),
    ]
    x0, y0 = Inches(0.55), Inches(1.95)
    cw, ch = Inches(6.1), Inches(1.35)
    for i, (head, body, col) in enumerate(ideas):
        r = i // 2; c = i % 2
        x = x0 + c * (cw + Inches(0.2))
        y = y0 + r * (ch + Inches(0.18))
        add_round_rect(s, x, y, cw, ch, LIGHT)
        add_rect(s, x, y, cw, Inches(0.08), col)
        add_text(s, x + Inches(0.25), y + Inches(0.2),
                 cw - Inches(0.5), Inches(0.45),
                 head, size=16, bold=True, color=col)
        add_text(s, x + Inches(0.25), y + Inches(0.7),
                 cw - Inches(0.5), ch - Inches(0.8),
                 body, size=12, color=DARK)


# ---------------------------------------------------------------------------
# Slide 20: Thanks / Q&A
# ---------------------------------------------------------------------------
@slide
def s_thanks(s):
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, PRIMARY)
    add_rect(s, 0, Inches(3.4), SLIDE_W, Inches(0.08), ACCENT)
    add_text(s, Inches(0.5), Inches(1.9), SLIDE_W - Inches(1), Inches(1.5),
             "Thank you", size=80, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(3.6), SLIDE_W - Inches(1), Inches(0.8),
             "Questions?", size=32, color=ACCENT, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(4.6), SLIDE_W - Inches(1), Inches(0.6),
             "COMP258  ·  Student Success Prediction",
             size=18, color=RGBColor(0xCB, 0xD5, 0xE1), align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(5.2), SLIDE_W - Inches(1), Inches(0.5),
             "Group:  [YourGroupName]  ·  Week 14",
             size=14, color=RGBColor(0xCB, 0xD5, 0xE1), align=PP_ALIGN.CENTER)


# ---------------------------------------------------------------------------
# Render
# ---------------------------------------------------------------------------
total = len(SLIDES)
for i, fn in enumerate(SLIDES, start=1):
    s = prs.slides.add_slide(BLANK)
    # light background
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    fn(s)
    # footer on all non-title/non-thanks slides
    if fn.__name__ not in ("s_title", "s_thanks"):
        add_page_footer(s, i, total)

prs.save(OUT_PPTX)
print(f"Wrote {OUT_PPTX}")
print(f"Slides: {total}")
